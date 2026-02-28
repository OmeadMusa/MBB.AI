#!/usr/bin/env python3
"""
MBB-style PowerPoint deck generator using python-pptx.

Usage (CLI):
    python generate_pptx.py --output FILE --title TITLE [--subtitle TEXT] [--date TEXT]

Usage (Python):
    from generate_pptx import MBBDeckBuilder
    deck = MBBDeckBuilder("Title", subtitle="Client", date="February 2026")
    deck.add_title_slide()
    deck.add_exec_summary(governing_thought, key_points)
    deck.add_content_slide(action_title, body_text, source)
    deck.save("output.pptx")
"""

from __future__ import annotations

import argparse
import os
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── MBB Color Constants ─────────────────────────────────────────────────────

NAVY = RGBColor(0x00, 0x3A, 0x70)
STEEL_BLUE = RGBColor(0x44, 0x72, 0xC4)
TEAL = RGBColor(0x00, 0xB0, 0xF0)
GREEN = RGBColor(0x00, 0xB0, 0x50)
RED = RGBColor(0xFF, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x80, 0x80, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

# ── Layout Constants (16:9 widescreen) ───────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_TOP = Inches(0.3)
TITLE_HEIGHT = Inches(0.9)
CONTENT_TOP = Inches(1.4)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
SOURCE_TOP = Inches(7.0)
SOURCE_HEIGHT = Inches(0.3)
PAGE_NUM_LEFT = Inches(12.5)

FONT_NAME = "Calibri"


class MBBDeckBuilder:
    """Creates MBB-standard PowerPoint decks following the Pyramid Principle."""

    def __init__(self, title: str, subtitle: str = "", date: str = None,
                 color_scheme: dict = None):
        self.title = title
        self.subtitle = subtitle
        self.date = date or datetime.now().strftime("%B %Y")
        self.slide_count = 0

        # Allow custom color overrides
        self.colors = {
            "primary": NAVY,
            "secondary": STEEL_BLUE,
            "accent": TEAL,
            "positive": GREEN,
            "negative": RED,
            "text": DARK_GRAY,
            "white": WHITE,
        }
        if color_scheme:
            self.colors.update(color_scheme)

        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def _get_blank_layout(self):
        """Get the blank slide layout for full control."""
        return self.prs.slide_layouts[6]

    def _add_slide(self):
        """Add a blank slide and increment counter."""
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self.slide_count += 1
        return slide

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=12, bold=False, italic=False, color=None,
                     alignment=PP_ALIGN.LEFT, font_name=FONT_NAME):
        """Add a formatted textbox to a slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color or self.colors["text"]
        p.font.name = font_name
        p.alignment = alignment

        return txBox

    def _add_action_title(self, slide, title_text):
        """Add an action title to a slide (the key message as an assertion)."""
        self._add_textbox(
            slide, MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
            title_text, font_size=22, bold=True, color=self.colors["primary"]
        )

        # Add thin separator line below title
        left = MARGIN
        top = Inches(1.2)
        width = CONTENT_WIDTH
        line = slide.shapes.add_connector(1, left, top, left + width, top)  # MSO_CONNECTOR.STRAIGHT = 1
        line.line.color.rgb = self.colors["secondary"]
        line.line.width = Pt(1)

    def _add_source_line(self, slide, source_text):
        """Add source citation at bottom of slide."""
        self._add_textbox(
            slide, MARGIN, SOURCE_TOP, Inches(10), SOURCE_HEIGHT,
            f"Source: {source_text}", font_size=8, italic=True,
            color=LIGHT_GRAY
        )

    def _add_page_number(self, slide):
        """Add page number to bottom right."""
        self._add_textbox(
            slide, PAGE_NUM_LEFT, SOURCE_TOP, Inches(0.5), SOURCE_HEIGHT,
            str(self.slide_count), font_size=8, color=LIGHT_GRAY,
            alignment=PP_ALIGN.RIGHT
        )

    def _add_bullets(self, slide, items, left, top, width, height,
                     font_size=12, color=None, spacing=Pt(6)):
        """Add a bulleted list to a slide."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color or self.colors["text"]
            p.font.name = FONT_NAME
            p.space_after = spacing
            p.level = 0

            # Add bullet
            pPr = p._pPr
            if pPr is None:
                from pptx.oxml.ns import qn
                pPr = p._p.get_or_add_pPr()
            from pptx.oxml.ns import qn
            buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
            # Remove existing bullet elements
            for child in list(pPr):
                if "buChar" in child.tag or "buNone" in child.tag:
                    pPr.remove(child)
            pPr.append(buChar)

        return txBox

    # ── Slide Types ──────────────────────────────────────────────────────────

    def add_title_slide(self):
        """Add the title slide with dark navy background."""
        slide = self._add_slide()

        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["primary"]

        # Title
        self._add_textbox(
            slide, Inches(1), Inches(2.2), Inches(11.333), Inches(1.5),
            self.title, font_size=36, bold=True, color=self.colors["white"],
            alignment=PP_ALIGN.CENTER
        )

        # Subtitle
        if self.subtitle:
            self._add_textbox(
                slide, Inches(1), Inches(3.8), Inches(11.333), Inches(0.8),
                self.subtitle, font_size=20, color=self.colors["white"],
                alignment=PP_ALIGN.CENTER
            )

        # Date
        self._add_textbox(
            slide, Inches(1), Inches(5.0), Inches(11.333), Inches(0.5),
            self.date, font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB),
            alignment=PP_ALIGN.CENTER
        )

        # Confidential watermark
        self._add_textbox(
            slide, Inches(1), Inches(6.5), Inches(11.333), Inches(0.4),
            "CONFIDENTIAL", font_size=10, color=RGBColor(0x88, 0x88, 0x88),
            alignment=PP_ALIGN.CENTER
        )

    def add_exec_summary(self, governing_thought: str, key_points: list[str]):
        """
        Add an executive summary slide.

        Args:
            governing_thought: The single most important message (becomes the title)
            key_points: 3-5 supporting points
        """
        slide = self._add_slide()
        self._add_action_title(slide, governing_thought)

        self._add_bullets(
            slide, key_points,
            left=Inches(0.8), top=CONTENT_TOP, width=Inches(11.5), height=Inches(5.0),
            font_size=14
        )

        self._add_page_number(slide)

    def add_content_slide(self, action_title: str, body_text: str, source: str = ""):
        """
        Add a standard content slide with action title and body text.

        Args:
            action_title: The key message as an assertion (NOT a topic label)
            body_text: Supporting content
            source: Data source citation
        """
        slide = self._add_slide()
        self._add_action_title(slide, action_title)

        content_height = Inches(5.0) if not source else Inches(4.5)
        self._add_textbox(
            slide, MARGIN, CONTENT_TOP, CONTENT_WIDTH, content_height,
            body_text, font_size=13
        )

        if source:
            self._add_source_line(slide, source)
        self._add_page_number(slide)

    def add_two_column_slide(self, action_title: str, left_content: str,
                              right_content: str, source: str = ""):
        """Add a two-column layout slide."""
        slide = self._add_slide()
        self._add_action_title(slide, action_title)

        col_width = Inches(5.8)
        content_height = Inches(4.5)

        # Left column
        self._add_textbox(
            slide, MARGIN, CONTENT_TOP, col_width, content_height,
            left_content, font_size=12
        )

        # Right column
        self._add_textbox(
            slide, Inches(7.0), CONTENT_TOP, col_width, content_height,
            right_content, font_size=12
        )

        if source:
            self._add_source_line(slide, source)
        self._add_page_number(slide)

    def add_chart_slide(self, action_title: str, chart_image_path: str,
                        source: str = ""):
        """
        Add a slide with an embedded chart image.

        Args:
            action_title: Key message as assertion
            chart_image_path: Path to chart image file (PNG/JPG)
            source: Data source citation
        """
        slide = self._add_slide()
        self._add_action_title(slide, action_title)

        # Center the chart image in the content area
        img_width = Inches(10)
        img_height = Inches(5.0)
        img_left = (SLIDE_WIDTH - img_width) // 2
        slide.shapes.add_picture(
            chart_image_path, img_left, CONTENT_TOP, img_width, img_height
        )

        if source:
            self._add_source_line(slide, source)
        self._add_page_number(slide)

    def add_table_slide(self, action_title: str, headers: list[str],
                        rows: list[list], source: str = ""):
        """
        Add a slide with a formatted data table.

        Args:
            action_title: Key message as assertion
            headers: Column headers
            rows: List of row data (each row is a list)
            source: Data source citation
        """
        slide = self._add_slide()
        self._add_action_title(slide, action_title)

        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        table_width = CONTENT_WIDTH
        table_height = Inches(min(4.5, 0.4 * num_rows + 0.5))
        table_left = MARGIN

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, table_left, CONTENT_TOP, table_width, table_height
        )
        table = table_shape.table

        # Style header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.colors["white"]
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors["primary"]

        # Style data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = self.colors["text"]
                    paragraph.font.name = FONT_NAME
                    paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

                # Alternate row shading
                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

        if source:
            self._add_source_line(slide, source)
        self._add_page_number(slide)

    def add_recommendation_slide(self, recommendations: list[dict]):
        """
        Add a recommendations slide.

        Args:
            recommendations: List of dicts with keys: what, why, impact, timeline
        """
        slide = self._add_slide()
        self._add_action_title(slide, "Recommendations")

        num_rows = len(recommendations) + 1
        headers = ["#", "Recommendation", "Rationale", "Expected Impact", "Timeline"]
        num_cols = len(headers)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols, MARGIN, CONTENT_TOP,
            CONTENT_WIDTH, Inches(min(4.5, 0.5 * num_rows + 0.5))
        )
        table = table_shape.table

        # Set column widths
        col_widths = [Inches(0.5), Inches(3.5), Inches(3.5), Inches(2.8), Inches(2.0)]
        for i, width in enumerate(col_widths):
            table.columns[i].width = width

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.colors["white"]
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors["primary"]

        # Data rows
        for row_idx, rec in enumerate(recommendations):
            values = [
                str(row_idx + 1),
                rec.get("what", ""),
                rec.get("why", ""),
                rec.get("impact", ""),
                rec.get("timeline", ""),
            ]
            for col_idx, value in enumerate(values):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = value
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = self.colors["text"]
                    paragraph.font.name = FONT_NAME
                    paragraph.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT

                if row_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

        self._add_page_number(slide)

    def add_next_steps_slide(self, immediate: list[str] = None,
                              short_term: list[str] = None,
                              medium_term: list[str] = None):
        """
        Add a next steps slide with phased actions.

        Args:
            immediate: 0-2 week actions
            short_term: 1-3 month actions
            medium_term: 3-12 month actions
        """
        slide = self._add_slide()
        self._add_action_title(slide, "Next Steps")

        phases = [
            ("Immediate (0-2 weeks)", immediate or []),
            ("Short-term (1-3 months)", short_term or []),
            ("Medium-term (3-12 months)", medium_term or []),
        ]

        col_width = Inches(3.8)
        for i, (phase_title, items) in enumerate(phases):
            left = Inches(0.5 + i * 4.2)

            # Phase header
            self._add_textbox(
                slide, left, CONTENT_TOP, col_width, Inches(0.5),
                phase_title, font_size=14, bold=True, color=self.colors["secondary"],
                alignment=PP_ALIGN.LEFT
            )

            # Phase items
            if items:
                self._add_bullets(
                    slide, items,
                    left=left, top=Inches(2.0), width=col_width, height=Inches(4.5),
                    font_size=11
                )

        self._add_page_number(slide)

    def add_section_divider(self, section_title: str):
        """Add a section break slide."""
        slide = self._add_slide()

        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["secondary"]

        self._add_textbox(
            slide, Inches(1), Inches(3.0), Inches(11.333), Inches(1.2),
            section_title, font_size=32, bold=True, color=self.colors["white"],
            alignment=PP_ALIGN.CENTER
        )

        self._add_page_number(slide)

    def add_appendix_divider(self):
        """Add an appendix section break."""
        self.add_section_divider("Appendix")

    # ── Save ─────────────────────────────────────────────────────────────────

    def save(self, filepath: str):
        """Save the presentation."""
        os.makedirs(os.path.dirname(filepath) if os.path.dirname(filepath) else ".", exist_ok=True)
        self.prs.save(filepath)
        print(f"Presentation saved: {filepath} ({self.slide_count} slides)")
        return filepath


# ── CLI Interface ────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Generate MBB-style PowerPoint deck")
    parser.add_argument("--output", required=True, help="Output file path")
    parser.add_argument("--title", required=True, help="Deck title")
    parser.add_argument("--subtitle", default="", help="Subtitle / client name")
    parser.add_argument("--date", default=None, help="Date string")
    args = parser.parse_args()

    deck = MBBDeckBuilder(args.title, subtitle=args.subtitle, date=args.date)
    deck.add_title_slide()
    deck.save(args.output)


if __name__ == "__main__":
    main()
